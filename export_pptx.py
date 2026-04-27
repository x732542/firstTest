"""
export_pptx.py
把 presentation.html 每張投影片截圖並合成 PPTX
執行方式: python export_pptx.py
"""
import asyncio, os, pathlib
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

HTML_FILE = pathlib.Path(__file__).parent / "presentation.html"
OUT_DIR   = pathlib.Path(__file__).parent / "_slides_tmp"
OUT_PPTX  = pathlib.Path(__file__).parent / "presentation.pptx"

SLIDE_W = 1200   # px — 對應 CSS 中 deck-shell 的寬度
SLIDE_H = 680    # px — 對應 CSS 中 deck-shell 的高度
TOTAL   = 10     # 投影片總數

async def capture_slides():
    OUT_DIR.mkdir(exist_ok=True)
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(viewport={"width": SLIDE_W, "height": SLIDE_H})
        await page.goto(HTML_FILE.as_uri(), wait_until="networkidle")

        # 先等字型與 Bootstrap Icons 載入
        await page.wait_for_timeout(1500)

        paths = []
        for i in range(TOTAL):
            # 呼叫頁面的 goTo 函數切換投影片
            await page.evaluate(f"goTo({i})")
            await page.wait_for_timeout(600)   # 等動畫結束

            # 只截投影片容器元素
            el = page.locator("#deck")
            out = OUT_DIR / f"slide_{i:02d}.png"
            await el.screenshot(path=str(out))
            paths.append(out)
            print(f"  截圖 slide {i+1}/{TOTAL}: {out.name}")

        await browser.close()
    return paths

def build_pptx(img_paths):
    prs = Presentation()
    # 設定投影片尺寸為 16:9 (33.87 cm × 19.05 cm)
    prs.slide_width  = Inches(SLIDE_W / 96)   # 96 dpi
    prs.slide_height = Inches(SLIDE_H / 96)

    blank_layout = prs.slide_layouts[6]   # 全空白版面

    for i, img_path in enumerate(img_paths):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path),
            left=0, top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
        print(f"  加入 PPTX 第 {i+1} 頁")

    prs.save(str(OUT_PPTX))
    print(f"\n✅ 已儲存：{OUT_PPTX}")

async def main():
    print("▶ 開始截圖...")
    paths = await capture_slides()
    print("\n▶ 合成 PPTX...")
    build_pptx(paths)

    # 清理暫存截圖
    for p in paths:
        p.unlink()
    OUT_DIR.rmdir()

asyncio.run(main())
